"""Document extraction with structure-first local processing and OCR fallback.

Local mode keeps native text where it exists and uses Tesseract only for empty or
low-text pages. Azure mode can use Document Intelligence prebuilt-layout once
approved/configured, without changing the ingestion API.
"""
from __future__ import annotations

import io
import json
import logging
import shutil
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

from ..config import settings

logger = logging.getLogger(__name__)


@dataclass
class ExtractionResult:
    text: str
    extractor: str
    page_count: int = 0
    warnings: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)

    def as_metadata(self) -> Dict[str, Any]:
        value = asdict(self)
        value.pop("text", None)
        return value


def extract_document(path: str, extension: str) -> ExtractionResult:
    """Extract readable, source-labelled text and return processing metadata."""
    extension = extension.lower().lstrip(".")
    processor = settings.DOCUMENT_PROCESSOR.lower()
    if processor in {"azure", "azure_document_intelligence"}:
        return _extract_with_document_intelligence(path)
    if processor == "auto" and extension in {"pdf", "docx", "xlsx", "pptx"}:
        try:
            return _extract_with_document_intelligence(path)
        except Exception as exc:
            logger.info("Document Intelligence unavailable; using local extractor: %s", type(exc).__name__)
            local = _extract_local(path, extension)
            local.warnings.append("Azure Document Intelligence was unavailable; processed locally.")
            return local
    return _extract_local(path, extension)


def extract_text(path: str, extension: str) -> str:
    """Compatibility wrapper for callers that need only text."""
    return extract_document(path, extension).text


def _extract_local(path: str, extension: str) -> ExtractionResult:
    if extension == "txt":
        return ExtractionResult(Path(path).read_text(encoding="utf-8", errors="replace"), "native_text")
    if extension == "json":
        content = json.loads(Path(path).read_text(encoding="utf-8"))
        return ExtractionResult(json.dumps(content, indent=2, ensure_ascii=False), "json")
    if extension == "docx":
        return _extract_docx(path)
    if extension == "pptx":
        return _extract_pptx(path)
    if extension == "xlsx":
        return _extract_xlsx(path)
    if extension == "pdf":
        return _extract_pdf(path)
    if extension in {"png", "jpg", "jpeg", "tif", "tiff"}:
        return _extract_image(path)
    raise ValueError(f"No local extractor configured for .{extension}")


def _extract_docx(path: str) -> ExtractionResult:
    from docx import Document

    document = Document(path)
    parts = ["=== Word document ==="]
    for index, paragraph in enumerate(document.paragraphs, 1):
        if paragraph.text.strip():
            parts.append(f"[Paragraph {index}] {paragraph.text.strip()}")
    for table_number, table in enumerate(document.tables, 1):
        rows = [" | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells) for row in table.rows]
        parts.append(f"=== Table {table_number} ===\n" + "\n".join(rows))
    return ExtractionResult("\n\n".join(parts), "python_docx", metadata={"tables": len(document.tables)})


def _extract_pptx(path: str) -> ExtractionResult:
    from pptx import Presentation

    presentation = Presentation(path)
    slides: List[str] = []
    warnings: List[str] = []
    for number, slide in enumerate(presentation.slides, 1):
        parts = [f"=== Slide {number} ==="]
        image_count = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                table_rows = [" | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells) for row in shape.table.rows]
                parts.append("[Table]\n" + "\n".join(table_rows))
            if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE without extra dependency
                image_count += 1
                ocr_text = _ocr_image_bytes(getattr(shape.image, "blob", b"")) if settings.OCR_ENABLED else ""
                if ocr_text:
                    parts.append(f"[Image OCR {image_count}] {ocr_text}")
        if len(parts) == 1 and settings.OCR_ENABLED:
            warnings.append(f"Slide {number} contains no extractable text; image OCR may be incomplete.")
        slides.append("\n".join(parts))
    return ExtractionResult("\n\n".join(slides), "python_pptx", page_count=len(presentation.slides), warnings=warnings)


def _extract_xlsx(path: str) -> ExtractionResult:
    import pandas as pd

    workbook = pd.ExcelFile(path)
    sheets: List[str] = []
    for sheet_name in workbook.sheet_names:
        frame = pd.read_excel(workbook, sheet_name=sheet_name, dtype=str).fillna("")
        # CSV is intentionally retained: rows/columns remain legible to the model.
        sheets.append(f"=== Sheet: {sheet_name} ===\n{frame.to_csv(index=False)}")
    return ExtractionResult("\n\n".join(sheets), "pandas_openpyxl", metadata={"sheets": workbook.sheet_names})


def _extract_pdf(path: str) -> ExtractionResult:
    from PyPDF2 import PdfReader

    reader = PdfReader(path)
    native_pages = [(page.extract_text() or "").strip() for page in reader.pages]
    low_text_pages = [index + 1 for index, text in enumerate(native_pages) if len(text) < settings.OCR_MIN_TEXT_PER_PAGE]
    ocr_pages: Dict[int, str] = {}
    warnings: List[str] = []
    if low_text_pages and settings.OCR_ENABLED:
        ocr_pages, ocr_warning = _ocr_pdf_pages(path, low_text_pages)
        if ocr_warning:
            warnings.append(ocr_warning)
    elif low_text_pages:
        warnings.append(f"Pages {low_text_pages} have little or no native text; OCR is disabled.")

    parts: List[str] = []
    for number, native_text in enumerate(native_pages, 1):
        page_text = native_text or ocr_pages.get(number, "")
        if not page_text:
            parts.append(f"--- Page {number} ---\n[No readable text extracted from this page]")
        else:
            source = "OCR" if not native_text and number in ocr_pages else "Native text"
            parts.append(f"--- Page {number} ({source}) ---\n{page_text}")
    extractor = "pypdf2+tesseract" if ocr_pages else "pypdf2"
    return ExtractionResult("\n\n".join(parts), extractor, page_count=len(reader.pages), warnings=warnings, metadata={"ocr_pages": sorted(ocr_pages)})


def _extract_image(path: str) -> ExtractionResult:
    text = _ocr_image_bytes(Path(path).read_bytes())
    if not text:
        raise ValueError(_ocr_install_message("The image could not be read"))
    return ExtractionResult(f"=== Image OCR ===\n{text}", "tesseract", page_count=1)


def _ocr_pdf_pages(path: str, page_numbers: List[int]) -> tuple[Dict[int, str], Optional[str]]:
    if not shutil.which("tesseract"):
        return {}, _ocr_install_message("Tesseract is not installed")
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError:
        return {}, "Install pdf2image and pytesseract to enable local OCR."
    if not shutil.which("pdftoppm") and not shutil.which("pdftocairo"):
        return {}, "Poppler is not installed; run `brew install poppler` to enable PDF-to-image OCR."
    try:
        images = convert_from_path(path, dpi=settings.OCR_DPI)
        result = {}
        for number in page_numbers:
            text = pytesseract.image_to_string(images[number - 1]).strip()
            if text:
                result[number] = text
        return result, None
    except Exception as exc:
        logger.warning("Tesseract PDF OCR failed: %s", type(exc).__name__)
        return {}, "Local OCR could not process one or more pages. Upload a clearer scan or configure Azure Document Intelligence."


def _ocr_image_bytes(blob: bytes) -> str:
    if not blob or not settings.OCR_ENABLED or not shutil.which("tesseract"):
        return ""
    try:
        from PIL import Image
        import pytesseract

        return pytesseract.image_to_string(Image.open(io.BytesIO(blob))).strip()
    except Exception:
        return ""


def _ocr_install_message(prefix: str) -> str:
    return f"{prefix}. Local OCR requires Tesseract (`brew install tesseract`) and Poppler (`brew install poppler`) on macOS."


def _extract_with_document_intelligence(path: str) -> ExtractionResult:
    if not all([settings.DOCUMENT_INTELLIGENCE_ENDPOINT, settings.DOCUMENT_INTELLIGENCE_API_KEY]):
        raise ValueError("Azure Document Intelligence is not configured")
    try:
        from azure.ai.documentintelligence import DocumentIntelligenceClient
        from azure.ai.documentintelligence.models import ContentFormat
        from azure.core.credentials import AzureKeyCredential
    except ImportError as exc:
        raise ValueError("Install azure-ai-documentintelligence to enable Azure document processing") from exc
    client = DocumentIntelligenceClient(
        endpoint=settings.DOCUMENT_INTELLIGENCE_ENDPOINT,
        credential=AzureKeyCredential(settings.DOCUMENT_INTELLIGENCE_API_KEY),
    )
    with open(path, "rb") as document_file:
        poller = client.begin_analyze_document(
            "prebuilt-layout",
            body=document_file,
            output_content_format=ContentFormat.MARKDOWN,
        )
    result = poller.result()
    content = (result.content or "").strip()
    if not content:
        raise ValueError("Azure Document Intelligence returned no readable content")
    tables = len(result.tables or [])
    pages = len(result.pages or [])
    return ExtractionResult(
        content,
        "azure_document_intelligence_prebuilt_layout",
        page_count=pages,
        metadata={"tables": tables, "model_id": getattr(result, "model_id", "prebuilt-layout")},
    )

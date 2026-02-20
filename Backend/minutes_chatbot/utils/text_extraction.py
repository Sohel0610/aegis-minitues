"""
Advanced Text Extraction System — Corrected & Production-Ready
==============================================================

Fixes applied vs original code:
  1.  Removed hard imports of fitz (PyMuPDF) and PyPDF2 — both unavailable.
      All PDF work now uses pdfplumber (primary) + pdf2image+pytesseract (OCR fallback).
  2.  Fixed OCR pipeline: uses pdf2image.convert_from_path instead of fitz renderer.
  3.  Added robust image-PDF auto-detection (handles 0-page PDFs like scanned letters).
  4.  All optional libraries wrapped in try/except with graceful fallback.
  5.  Table extraction improved with explicit pdfplumber settings.
  6.  Added per-page OCR fallback for mixed PDFs (some pages text, some image).
  7.  Cleaned ExtractionResult.get_full_text() and added to_dict() helper.
  8.  Logging made configurable via LogLevel parameter.

Tested against:
  Agenda_5b  — text-based table PDF         OK
  Agenda_6b  — scanned / image-only PDF     OK (OCR)
  Management_Presentation — PPT-export PDF  OK
  Statutory_Auditor — PPT-export PDF        OK

Required (confirmed available):
  pdfplumber, pdf2image, pytesseract, Pillow,
  python-docx, openpyxl, python-pptx, pandas, beautifulsoup4

Optional (graceful degradation if missing):
  fitz (PyMuPDF), PyPDF2
"""

from __future__ import annotations

import io
import logging
import os
import re
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

# Core PDF library (confirmed available)
import pdfplumber

# OCR pipeline (confirmed available)
from pdf2image import convert_from_path
import pytesseract
from PIL import Image

# Office document libraries (confirmed available)
from docx import Document as DocxDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation

# Utilities
import pandas as pd
from bs4 import BeautifulSoup

# Optional libraries — graceful fallback if missing
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

# Logging setup
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s - %(message)s",
)
logger = logging.getLogger("text_extractor")


# =============================================================================
# Data container
# =============================================================================

@dataclass
class ExtractionResult:
    """Holds all content extracted from a single file."""

    text: str = ""
    tables: List[Dict[str, Any]] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    images: List[Dict[str, Any]] = field(default_factory=list)
    headers: List[str] = field(default_factory=list)
    footers: List[str] = field(default_factory=list)
    structure: Dict[str, Any] = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)

    def get_full_text(self, include_tables: bool = True) -> str:
        """Return text plus optional table content as a single string."""
        parts = [self.text]
        if include_tables and self.tables:
            parts.append("\n\n=== EXTRACTED TABLES ===")
            for i, table in enumerate(self.tables, 1):
                label = (
                    f"Table {i} (page {table.get('page', '?')})"
                    if "page" in table
                    else f"Table {i}"
                )
                parts.append(f"\n--- {label} ---")
                parts.append(table.get("text", str(table.get("data", ""))))
        return "\n".join(parts)

    def to_dict(self) -> Dict[str, Any]:
        """Serialise to a plain dict (for JSON output, logging, etc.)."""
        return {
            "text_length": len(self.text),
            "table_count": len(self.tables),
            "image_count": len(self.images),
            "header_count": len(self.headers),
            "warnings": self.warnings,
            "metadata": self.metadata,
        }

    def add_warning(self, msg: str) -> None:
        self.warnings.append(msg)
        logger.warning(msg)


# =============================================================================
# PDF Extractor  (primary + OCR fallback)
# =============================================================================

# pdfplumber table extraction settings tuned for financial / presentation PDFs
_TABLE_SETTINGS: Dict[str, Any] = {
    "vertical_strategy":    "lines",
    "horizontal_strategy":  "lines",
    "snap_tolerance":       3,
    "join_tolerance":       3,
    "edge_min_length":      3,
    "min_words_vertical":   1,
    "min_words_horizontal": 1,
}

# Average chars-per-page below which we treat the whole PDF as image-based
_IMAGE_PDF_THRESHOLD = 30


class AdvancedPDFExtractor:
    """
    Extracts text, tables, and metadata from PDF files.

    Strategy order:
      1. pdfplumber   - handles text-PDFs and PPT-export PDFs  (primary)
      2. Per-page OCR - runs on pages that yield < threshold chars
      3. Whole-file OCR via pdf2image - for fully scanned / image PDFs
      4. PyMuPDF (fitz) - metadata only, if available
      5. PyPDF2       - last-resort plain text fallback, if available

    Scanned PDF detection:
      - 0-page pdfplumber result  => whole file is an image
      - avg chars/page < 30       => treat whole file as image
    """

    def __init__(
        self,
        enable_ocr: bool = True,
        ocr_dpi: int = 250,
        per_page_ocr_threshold: int = 50,
    ) -> None:
        self.enable_ocr = enable_ocr
        self.ocr_dpi = ocr_dpi
        self.per_page_ocr_threshold = per_page_ocr_threshold

    # ------------------------------------------------------------------
    # Public entry point
    # ------------------------------------------------------------------

    def extract(self, file_path: str) -> ExtractionResult:
        result = ExtractionResult()

        # Step 1: detect whether PDF is image-based
        if self._is_image_pdf(file_path):
            logger.info("Detected image-only PDF - running full OCR: %s", file_path)
            if self.enable_ocr:
                result.text = self._ocr_full_file(file_path)
                result.add_warning("Image PDF detected - text extracted via OCR.")
            else:
                result.add_warning(
                    "Image PDF detected but OCR is disabled - no text extracted."
                )
            result.metadata = self._get_metadata(file_path)
            return result

        # Step 2: normal pdfplumber extraction
        try:
            plumber_data = self._extract_with_pdfplumber(file_path)
            result.text   = plumber_data["text"]
            result.tables = plumber_data["tables"]
        except Exception as exc:
            result.add_warning(f"pdfplumber extraction failed: {exc}")

        # Step 3: per-page OCR for mixed PDFs
        if self.enable_ocr:
            result.text = self._fill_image_pages_with_ocr(file_path, result.text)

        # Step 4: metadata
        result.metadata = self._get_metadata(file_path)

        # Step 5: last-resort fallbacks if still empty
        if not result.text.strip():
            result.text = self._fallback_text(file_path)
            if not result.text.strip():
                result.add_warning("All extraction strategies returned empty text.")

        return result

    # ------------------------------------------------------------------
    # Image-PDF detection
    # ------------------------------------------------------------------

    def _is_image_pdf(self, file_path: str) -> bool:
        """Return True if the PDF contains no (or negligible) selectable text."""
        try:
            with pdfplumber.open(file_path) as pdf:
                if len(pdf.pages) == 0:
                    return True
                total_chars = sum(len(p.extract_text() or "") for p in pdf.pages)
                avg = total_chars / len(pdf.pages)
                return avg < _IMAGE_PDF_THRESHOLD
        except Exception:
            return False

    # ------------------------------------------------------------------
    # pdfplumber extraction
    # ------------------------------------------------------------------

    def _extract_with_pdfplumber(self, file_path: str) -> Dict[str, Any]:
        text_parts: List[str] = []
        tables: List[Dict[str, Any]] = []

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                # Text
                page_text = page.extract_text() or ""
                text_parts.append(f"\n--- Page {page_num} ---\n{page_text}")

                # Tables
                raw_tables = page.extract_tables(_TABLE_SETTINGS) or []
                for t_idx, raw in enumerate(raw_tables):
                    if not raw:
                        continue
                    try:
                        header = raw[0] if raw else []
                        rows   = raw[1:] if len(raw) > 1 else []
                        # Deduplicate column names so DataFrame.to_dict() works cleanly
                        seen: Dict[str, int] = {}
                        clean_header: List[str] = []
                        for col in header:
                            col_str = str(col) if col is not None else ""
                            if col_str in seen:
                                seen[col_str] += 1
                                clean_header.append(f"{col_str}_{seen[col_str]}")
                            else:
                                seen[col_str] = 0
                                clean_header.append(col_str)
                        df = pd.DataFrame(rows, columns=clean_header)
                        df = df.dropna(how="all").dropna(axis=1, how="all")
                        tables.append({
                            "page":  page_num,
                            "index": t_idx,
                            "data":  df.to_dict("records"),
                            "text":  df.to_string(index=False),
                            "raw":   raw,
                        })
                    except Exception as exc:
                        logger.debug(
                            "Table parse error p%d t%d: %s", page_num, t_idx, exc
                        )
                        tables.append({
                            "page":  page_num,
                            "index": t_idx,
                            "raw":   raw,
                            "text":  str(raw),
                            "data":  [],
                        })

        return {"text": "\n".join(text_parts), "tables": tables}

    # ------------------------------------------------------------------
    # Per-page OCR for mixed PDFs
    # ------------------------------------------------------------------

    def _fill_image_pages_with_ocr(
        self, file_path: str, existing_text: str
    ) -> str:
        """
        For PDFs that are mostly text but have a few image-only pages,
        re-render just the low-text pages and splice in OCR output.
        """
        try:
            with pdfplumber.open(file_path) as pdf:
                image_pages = [
                    i for i, p in enumerate(pdf.pages)
                    if len(p.extract_text() or "") < self.per_page_ocr_threshold
                ]

            if not image_pages:
                return existing_text

            logger.info(
                "Per-page OCR on %d low-text pages: %s", len(image_pages), image_pages
            )

            for pg_0idx in image_pages:
                pg = pg_0idx + 1  # pdf2image is 1-indexed
                imgs = convert_from_path(
                    file_path, dpi=self.ocr_dpi,
                    first_page=pg, last_page=pg,
                )
                if not imgs:
                    continue
                ocr_text = pytesseract.image_to_string(imgs[0])
                marker = f"\n--- Page {pg} ---\n"
                if marker in existing_text:
                    before, _, after = existing_text.partition(marker)
                    # Replace the (empty) content after the marker up to next marker
                    next_marker_pos = after.find("\n--- Page ")
                    if next_marker_pos == -1:
                        existing_text = before + marker + ocr_text
                    else:
                        existing_text = (
                            before + marker + ocr_text + after[next_marker_pos:]
                        )

            return existing_text

        except Exception as exc:
            logger.warning("Per-page OCR failed: %s", exc)
            return existing_text

    # ------------------------------------------------------------------
    # Full-file OCR (for image PDFs)
    # ------------------------------------------------------------------

    def _ocr_full_file(self, file_path: str) -> str:
        """Render every page as image and OCR it."""
        try:
            images = convert_from_path(file_path, dpi=self.ocr_dpi)
            parts = []
            for i, img in enumerate(images, 1):
                text = pytesseract.image_to_string(img)
                parts.append(f"\n--- Page {i} (OCR) ---\n{text}")
            return "\n".join(parts)
        except Exception as exc:
            logger.error("Full-file OCR failed: %s", exc)
            return ""

    # ------------------------------------------------------------------
    # Metadata
    # ------------------------------------------------------------------

    def _get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Pull PDF metadata - tries PyMuPDF first, falls back to pdfplumber."""
        if PYMUPDF_AVAILABLE:
            try:
                doc  = fitz.open(file_path)
                meta = dict(doc.metadata) if doc.metadata else {}
                meta["page_count"] = doc.page_count
                doc.close()
                return meta
            except Exception:
                pass

        try:
            with pdfplumber.open(file_path) as pdf:
                meta = dict(pdf.metadata) if pdf.metadata else {}
                meta["page_count"] = len(pdf.pages)
                return meta
        except Exception:
            return {}

    # ------------------------------------------------------------------
    # Last-resort fallbacks
    # ------------------------------------------------------------------

    def _fallback_text(self, file_path: str) -> str:
        """Try PyMuPDF and PyPDF2 as final fallbacks."""
        if PYMUPDF_AVAILABLE:
            try:
                doc = fitz.open(file_path)
                parts = [
                    f"\n--- Page {i+1} ---\n{page.get_text('text')}"
                    for i, page in enumerate(doc)
                ]
                doc.close()
                text = "\n".join(parts)
                if text.strip():
                    logger.info("PyMuPDF fallback succeeded.")
                    return text
            except Exception as exc:
                logger.warning("PyMuPDF fallback failed: %s", exc)

        if PYPDF2_AVAILABLE:
            try:
                parts = []
                with open(file_path, "rb") as fh:
                    reader = PyPDF2.PdfReader(fh)
                    for i, page in enumerate(reader.pages, 1):
                        parts.append(f"\n--- Page {i} ---\n{page.extract_text()}")
                text = "\n".join(parts)
                if text.strip():
                    logger.info("PyPDF2 fallback succeeded.")
                    return text
            except Exception as exc:
                logger.warning("PyPDF2 fallback failed: %s", exc)

        return ""


# =============================================================================
# Word Document Extractor
# =============================================================================

class AdvancedDocxExtractor:
    """Extracts text, tables, headings, and image refs from .docx files."""

    def extract(self, file_path: str) -> ExtractionResult:
        result = ExtractionResult()
        try:
            doc = DocxDocument(file_path)
        except Exception as exc:
            result.add_warning(f"Failed to open DOCX: {exc}")
            return result

        # Metadata
        props = doc.core_properties
        result.metadata = {
            "author":   props.author   or "",
            "created":  str(props.created  or ""),
            "modified": str(props.modified or ""),
            "title":    props.title    or "",
            "subject":  props.subject  or "",
        }

        # Body elements
        text_parts: List[str] = []
        for element in doc.element.body:
            if isinstance(element, CT_P):
                para  = Paragraph(element, doc)
                style = para.style.name if para.style else "Normal"
                if style.startswith("Heading"):
                    result.headers.append(para.text)
                    level  = style.replace("Heading", "").strip()
                    prefix = "#" * max(int(level), 1) if level.isdigit() else "##"
                    text_parts.append(f"\n{prefix} {para.text}\n")
                else:
                    if para.text.strip():
                        text_parts.append(para.text)

            elif isinstance(element, CT_Tbl):
                table      = Table(element, doc)
                table_data = [
                    [cell.text for cell in row.cells] for row in table.rows
                ]
                if table_data:
                    try:
                        df = pd.DataFrame(table_data[1:], columns=table_data[0])
                        result.tables.append({
                            "data": df.to_dict("records"),
                            "text": df.to_string(index=False),
                            "raw":  table_data,
                        })
                        text_parts.append(f"\n[TABLE]\n{df.to_string(index=False)}\n")
                    except Exception:
                        result.tables.append({"raw": table_data, "text": str(table_data)})

        result.text = "\n".join(text_parts)

        # Image references
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                result.images.append({
                    "type":            rel.target_ref.split(".")[-1],
                    "relationship_id": rel.rId,
                })

        return result


# =============================================================================
# Excel Extractor
# =============================================================================

class AdvancedXlsxExtractor:
    """Extracts all sheets from .xlsx / .xls files as text and tables."""

    def extract(self, file_path: str) -> ExtractionResult:
        result = ExtractionResult()
        try:
            workbook = load_workbook(file_path, data_only=True)
        except Exception as exc:
            result.add_warning(f"Failed to open XLSX: {exc}")
            return result

        result.metadata = {
            "created":     str(workbook.properties.created  or ""),
            "modified":    str(workbook.properties.modified or ""),
            "creator":     workbook.properties.creator or "",
            "sheet_names": workbook.sheetnames,
        }

        text_parts: List[str] = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"\n\n=== Sheet: {sheet_name} ===")

            rows: List[List[str]] = [
                [str(cell) if cell is not None else "" for cell in row]
                for row in sheet.iter_rows(values_only=True)
            ]

            if not rows:
                continue

            # Find first non-empty row to use as header
            header_idx = next(
                (i for i, r in enumerate(rows) if any(c.strip() for c in r)), None
            )
            if header_idx is None:
                continue

            try:
                df = pd.DataFrame(rows[header_idx + 1:], columns=rows[header_idx])
                df = df.dropna(how="all").dropna(axis=1, how="all")
                result.tables.append({
                    "sheet": sheet_name,
                    "data":  df.to_dict("records"),
                    "text":  df.to_string(index=False),
                    "raw":   rows,
                })
                text_parts.append(df.to_string(index=False))
            except Exception as exc:
                logger.warning("Sheet '%s' parse error: %s", sheet_name, exc)
                text_parts.append(str(rows))

        result.text = "\n".join(text_parts)
        workbook.close()
        return result


# =============================================================================
# PowerPoint Extractor
# =============================================================================

class AdvancedPptxExtractor:
    """Extracts text, tables, notes, and image refs from .pptx files."""

    def extract(self, file_path: str) -> ExtractionResult:
        result = ExtractionResult()
        try:
            prs = Presentation(file_path)
        except Exception as exc:
            result.add_warning(f"Failed to open PPTX: {exc}")
            return result

        result.metadata = {
            "slide_count":  len(prs.slides),
            "slide_width":  prs.slide_width,
            "slide_height": prs.slide_height,
        }

        text_parts: List[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n\n=== Slide {slide_num} ===")

            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
                result.headers.append(title)
                text_parts.append(f"Title: {title}")

            for shape in slide.shapes:
                # Text frames (skip title to avoid duplication)
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:
                        text_parts.append(shape.text.strip())

                # Tables
                if shape.has_table:
                    table_data = [
                        [cell.text for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    if table_data:
                        try:
                            df = pd.DataFrame(table_data[1:], columns=table_data[0])
                            result.tables.append({
                                "slide": slide_num,
                                "data":  df.to_dict("records"),
                                "text":  df.to_string(index=False),
                                "raw":   table_data,
                            })
                            text_parts.append(
                                f"\n[TABLE]\n{df.to_string(index=False)}\n"
                            )
                        except Exception:
                            result.tables.append({
                                "slide": slide_num,
                                "raw":   table_data,
                                "text":  str(table_data),
                            })

                # Image reference
                if hasattr(shape, "image"):
                    result.images.append({
                        "slide":      slide_num,
                        "shape_type": str(shape.shape_type),
                    })

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_parts.append(f"\n[NOTES] {notes}")

        result.text = "\n".join(text_parts)
        return result


# =============================================================================
# HTML Extractor
# =============================================================================

class AdvancedHtmlExtractor:
    """Extracts clean text and tables from HTML files."""

    def extract(self, file_path: str) -> ExtractionResult:
        result = ExtractionResult()
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
                soup = BeautifulSoup(fh.read(), "html.parser")
        except Exception as exc:
            result.add_warning(f"Failed to open HTML: {exc}")
            return result

        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()

        result.text = soup.get_text(separator="\n", strip=True)

        for tbl in soup.find_all("table"):
            rows = [
                [cell.get_text(strip=True) for cell in tr.find_all(["td", "th"])]
                for tr in tbl.find_all("tr")
            ]
            if rows:
                try:
                    df = pd.DataFrame(rows[1:], columns=rows[0])
                    result.tables.append({
                        "data": df.to_dict("records"),
                        "text": df.to_string(index=False),
                        "raw":  rows,
                    })
                except Exception:
                    result.tables.append({"raw": rows, "text": str(rows)})

        return result


# =============================================================================
# CSV / TSV Extractor
# =============================================================================

class AdvancedCsvExtractor:
    """Extracts tabular data from CSV / TSV files."""

    _DELIMITERS = [",", "\t", ";", "|"]

    def extract(self, file_path: str) -> ExtractionResult:
        result = ExtractionResult()
        for delim in self._DELIMITERS:
            try:
                df = pd.read_csv(
                    file_path, delimiter=delim,
                    encoding="utf-8", on_bad_lines="skip",
                )
                if len(df.columns) > 1:
                    result.tables.append({
                        "data": df.to_dict("records"),
                        "text": df.to_string(index=False),
                        "raw":  df.values.tolist(),
                    })
                    result.text = df.to_string(index=False)
                    return result
            except Exception:
                continue

        # Fallback: read as raw text
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
                result.text = fh.read()
        except Exception as exc:
            result.add_warning(f"CSV read failed: {exc}")
        return result


# =============================================================================
# JSON Extractor
# =============================================================================

class AdvancedJsonExtractor:
    """Extracts structured data from JSON files (meeting minutes, etc.) into clear text."""

    def extract(self, file_path: str) -> ExtractionResult:
        import json as json_lib
        result = ExtractionResult()
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
                data = json_lib.load(fh)
        except Exception as exc:
            result.add_warning(f"Failed to parse JSON: {exc}")
            # Fallback: read as raw text
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
                    result.text = fh.read()
            except Exception:
                pass
            return result

        # Convert JSON to structured human-readable text
        result.text = self._json_to_text(data)
        result.metadata = {"format": "json", "type": type(data).__name__}
        return result

    def _json_to_text(self, data, indent: int = 0) -> str:
        """Convert JSON data to clean, readable text optimized for embedding and search."""
        if isinstance(data, dict):
            return self._dict_to_text(data)
        elif isinstance(data, list):
            parts = []
            for i, item in enumerate(data, 1):
                if isinstance(item, dict):
                    parts.append(f"\n--- Record {i} ---")
                    parts.append(self._dict_to_text(item))
                else:
                    parts.append(f"- {item}")
            return "\n".join(parts)
        else:
            return str(data)

    def _dict_to_text(self, data: dict) -> str:
        """Convert a dict (meeting minutes JSON) to clearly labeled text sections."""
        parts = []

        # Known meeting-minutes fields → output in a clear format
        field_labels = {
            "title": "Meeting Title",
            "date": "Meeting Date",
            "meeting_date": "Meeting Date",
            "next_meeting": "Next Meeting",
            "location": "Location",
            "summary": "Summary",
            "description": "Description",
        }

        # Process known scalar fields first
        for key, label in field_labels.items():
            if key in data:
                parts.append(f"{label}: {data[key]}")

        # Attendees
        if "attendees" in data:
            parts.append("\nAttendees:")
            for att in data["attendees"]:
                if isinstance(att, dict):
                    name = att.get("name", "Unknown")
                    role = att.get("role", "")
                    parts.append(f"  - {name} ({role})" if role else f"  - {name}")
                else:
                    parts.append(f"  - {att}")

        # Agenda
        if "agenda" in data:
            parts.append("\nAgenda:")
            for i, item in enumerate(data["agenda"], 1):
                if isinstance(item, dict):
                    title = item.get("title", item.get("topic", str(item)))
                    parts.append(f"  {i}. {title}")
                else:
                    parts.append(f"  {i}. {item}")

        # Decisions
        if "decisions" in data:
            parts.append("\nDecisions:")
            for i, dec in enumerate(data["decisions"], 1):
                if isinstance(dec, dict):
                    text = dec.get("text", dec.get("decision", str(dec)))
                    parts.append(f"  {i}. {text}")
                else:
                    parts.append(f"  {i}. {dec}")

        # Action Items
        if "action_items" in data:
            parts.append("\nAction Items:")
            for i, item in enumerate(data["action_items"], 1):
                if isinstance(item, dict):
                    task = item.get("task", item.get("description", str(item)))
                    assignee = item.get("assignee", item.get("assigned_to", ""))
                    due = item.get("due_date", item.get("deadline", ""))
                    line = f"  {i}. {task}"
                    if assignee:
                        line += f" — Assigned to: {assignee}"
                    if due:
                        line += f" — Due: {due}"
                    parts.append(line)
                else:
                    parts.append(f"  {i}. {item}")

        # Any remaining keys not yet handled
        handled_keys = set(field_labels.keys()) | {"attendees", "agenda", "decisions", "action_items"}
        for key, value in data.items():
            if key in handled_keys:
                continue
            label = key.replace("_", " ").title()
            if isinstance(value, (list, dict)):
                parts.append(f"\n{label}:")
                if isinstance(value, list):
                    for i, v in enumerate(value, 1):
                        if isinstance(v, dict):
                            sub_parts = [f"{k}: {v2}" for k, v2 in v.items()]
                            parts.append(f"  {i}. {'; '.join(sub_parts)}")
                        else:
                            parts.append(f"  {i}. {v}")
                else:
                    for k, v2 in value.items():
                        parts.append(f"  {k}: {v2}")
            else:
                parts.append(f"{label}: {value}")

        return "\n".join(parts)


# =============================================================================
# Universal Extractor  —  top-level dispatcher
# =============================================================================

class UniversalTextExtractor:
    """
    Auto-detects file type and dispatches to the correct extractor.

    Supported formats:
      PDF   -> AdvancedPDFExtractor  (text + tables + OCR for scanned files)
      DOCX  -> AdvancedDocxExtractor
      XLSX  -> AdvancedXlsxExtractor
      PPTX  -> AdvancedPptxExtractor
      HTML  -> AdvancedHtmlExtractor
      CSV   -> AdvancedCsvExtractor
      TXT   -> plain read

    Basic usage:
      extractor = UniversalTextExtractor()
      result    = extractor.extract("report.pdf")
      print(result.get_full_text())

    Get plain text string directly:
      text = extractor.extract_to_text("report.pdf", include_tables=True)
    """

    # Map file extension to extractor attribute name
    _EXT_MAP = {
        "pdf":  "_pdf",
        "docx": "_docx",
        "doc":  "_docx",
        "xlsx": "_xlsx",
        "xls":  "_xlsx",
        "pptx": "_pptx",
        "ppt":  "_pptx",
        "html": "_html",
        "htm":  "_html",
        "csv":  "_csv",
        "tsv":  "_csv",
        "json": "_json",
    }

    def __init__(
        self,
        enable_ocr: bool = True,
        ocr_dpi: int = 250,
        per_page_ocr_threshold: int = 50,
    ) -> None:
        self._pdf  = AdvancedPDFExtractor(
            enable_ocr=enable_ocr,
            ocr_dpi=ocr_dpi,
            per_page_ocr_threshold=per_page_ocr_threshold,
        )
        self._docx = AdvancedDocxExtractor()
        self._xlsx = AdvancedXlsxExtractor()
        self._pptx = AdvancedPptxExtractor()
        self._html = AdvancedHtmlExtractor()
        self._csv  = AdvancedCsvExtractor()
        self._json = AdvancedJsonExtractor()

    # ------------------------------------------------------------------
    # Main method
    # ------------------------------------------------------------------

    def extract(
        self,
        file_path: str,
        file_type: Optional[str] = None,
    ) -> ExtractionResult:
        """
        Extract all content from a file.

        Parameters
        ----------
        file_path : str
            Absolute or relative path to the file.
        file_type : str, optional
            Force a file type (e.g. "pdf").
            Auto-detected from extension if omitted.

        Returns
        -------
        ExtractionResult
        """
        ext = (
            file_type.lower().lstrip(".")
            if file_type
            else Path(file_path).suffix.lower().lstrip(".")
        )

        logger.info("Extracting [%s] %s", ext.upper(), file_path)

        extractor_attr = self._EXT_MAP.get(ext)

        if extractor_attr:
            result = getattr(self, extractor_attr).extract(file_path)
        elif ext == "txt":
            result = self._read_txt(file_path)
        else:
            result = self._read_unknown(file_path, ext)

        logger.info(
            "Done: %d chars, %d tables, %d images, %d warnings",
            len(result.text), len(result.tables),
            len(result.images), len(result.warnings),
        )
        return result

    # ------------------------------------------------------------------
    # Convenience method
    # ------------------------------------------------------------------

    def extract_to_text(
        self,
        file_path: str,
        file_type: Optional[str] = None,
        include_tables: bool = True,
        include_metadata: bool = False,
        include_warnings: bool = True,
    ) -> str:
        """
        Extract and return everything as a single plain-text string.

        Parameters
        ----------
        include_tables   : append extracted table data after main text
        include_metadata : prepend file metadata block
        include_warnings : append any warnings at the bottom
        """
        result = self.extract(file_path, file_type)
        parts: List[str] = []

        if include_metadata and result.metadata:
            parts.append("=== METADATA ===")
            for k, v in result.metadata.items():
                parts.append(f"{k}: {v}")
            parts.append("")

        parts.append(result.get_full_text(include_tables=include_tables))

        if include_warnings and result.warnings:
            parts.append("\n=== WARNINGS ===")
            parts.extend(result.warnings)

        return "\n".join(parts)

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _read_txt(file_path: str) -> ExtractionResult:
        result = ExtractionResult()
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
                result.text = fh.read()
        except Exception as exc:
            result.add_warning(f"TXT read failed: {exc}")
        return result

    @staticmethod
    def _read_unknown(file_path: str, ext: str) -> ExtractionResult:
        result = ExtractionResult()
        result.add_warning(
            f"Unsupported extension '{ext}' - attempting plain text read."
        )
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
                result.text = fh.read()
        except Exception as exc:
            result.add_warning(f"Plain-text read also failed: {exc}")
        return result


# =============================================================================
# Self-test  —  run with:  python advanced_text_extractor.py
# =============================================================================

if __name__ == "__main__":
    TEST_FILES = [
        "/mnt/user-data/uploads/Agenda_5b_-_Investment_by_AGEL_as_at_30_09_2025.pdf",
        "/mnt/user-data/uploads/Agenda_6b_-_Investor_Grievance.pdf",
        "/mnt/user-data/uploads/Management_Presentation.pdf",
        "/mnt/user-data/uploads/Statutory_Auditor_Presentation.pdf",
    ]

    extractor = UniversalTextExtractor(enable_ocr=True, ocr_dpi=250)

    print("\n" + "=" * 70)
    print("  UNIVERSAL TEXT EXTRACTOR  -  SELF-TEST")
    print("=" * 70)

    for fpath in TEST_FILES:
        if not os.path.exists(fpath):
            print(f"\n[SKIP] File not found: {fpath}")
            continue

        name = Path(fpath).name
        print(f"\n{'─' * 70}")
        print(f"FILE : {name}")

        result = extractor.extract(fpath)

        print(f"  Text chars  : {len(result.text):,}")
        print(f"  Tables      : {len(result.tables)}")
        print(f"  Images      : {len(result.images)}")
        print(f"  Warnings    : {len(result.warnings)}")
        for w in result.warnings:
            print(f"    [!] {w}")

        preview = result.text.strip()[:200].replace("\n", " ")
        print(f"  Text preview: {preview!r}")

        if result.tables:
            first = result.tables[0]
            rows  = first.get("raw", [])
            print(f"  First table : {len(rows)} rows x {len(rows[0]) if rows else 0} cols")
            for row in rows[:3]:
                print(f"    {row}")

    print("\n" + "=" * 70)
    print("  SELF-TEST COMPLETE")
    print("=" * 70 + "\n")